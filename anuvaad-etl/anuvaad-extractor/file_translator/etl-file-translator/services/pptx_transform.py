import copy
import os
import uuid

from anuvaad_auditor import log_info
from pptx import Presentation
from pptx.text.text import _Paragraph


import config
from errors.errors_exception import FileErrors
from services.service import common_obj


class PptxTransform(object):
    def __init__(self, input_filename, json_data):
        self.json_data = json_data
        self.file_name_without_ext = os.path.splitext(input_filename)[0]

        self.outer_struct = common_obj.outer_struct
        self.page_struct = common_obj.page_struct
        self.para_struct = common_obj.para_struct
        self.run_struct = common_obj.run_struct

    # reading content of pptx file
    def read_pptx_file(self, input_filename):
        input_pptx_filepath = common_obj.input_path(input_filename)
        input_pptx = Presentation(input_pptx_filepath)
        return input_pptx

    def generate_json_for_page(self, page_number):
        new_page_template = copy.deepcopy(self.page_struct)
        new_page_template['page_no'] = page_number
        return new_page_template

    '''generate_json_for_para accept para object and using that it tries to create a json structure'''

    def generate_json_for_para(self, para, file_id='', table='', cell='', row='', slide='', shape='', para_idx=''):
        new_para_template = copy.deepcopy(self.para_struct)
        runs = common_obj.get_runs(para, para_obj=True, file_type=config.TYPE_PPTX)
        new_para_template['text'] = common_obj.get_para_text(runs)
        new_para_template['block_id'] = common_obj.generate_id(file_id=file_id, table=table, cell=cell, row=row,
                                                               slide=slide, shape=shape, para=para_idx)
        return new_para_template

    def generate_json_for_run(self, run, file_id='', table='', cell='', row='', slide='', shape='', para_idx='',
                              run_idx=''):
        new_run_template = copy.deepcopy(self.run_struct)
        new_run_template['text'] = run.text
        new_run_template['block_id'] = common_obj.generate_id(file_id=file_id, table=table, cell=cell, row=row,
                                                              slide=slide, shape=shape, para=para_idx, run=run_idx)
        return new_run_template

    def generate_json_structure(self, document):
        base_json = copy.deepcopy(self.outer_struct)

        page_list = base_json['result']
        file_id = self.file_name_without_ext

        para_count = 0
        page_number = 1

        page_list.append(self.generate_json_for_page(page_number))

        # NEW LOGIC TO GET ALL PARAGRAPHS
        if config.PPTX_PARAGRAPH_GEN:
            try:
                for id_sld, slide in enumerate(document.slides):
                    try:
                        para_lis = slide._element.xpath('.//a:p')
                        for idx, para_obj in enumerate(para_lis):
                            try:
                                para = _Paragraph(para_obj, slide)
                                para_count += 1
                                json_para = self.generate_json_for_para(para=para,
                                                                        file_id=file_id,
                                                                        slide=str(id_sld),
                                                                        para_idx=str(idx))
                                for idr, run in enumerate(para.runs):
                                    json_run = self.generate_json_for_run(run=run,
                                                                          file_id=file_id,
                                                                          slide=str(id_sld),
                                                                          para_idx=str(idx),
                                                                          run_idx=str(idr))
                                    json_para['children'].append(json_run)
                                page_list[page_number - 1]['text_blocks'].append(json_para)
                            except Exception as e:
                                log_info(f"generate_json_structure :: PPTX Generate Json failed for Paragraph:{idx}, slide:{id_sld}, file:{file_id}"
                                         f"ERROR:{str(e)}", None)
                        page_number += 1
                        page_list.append(self.generate_json_for_page(page_number))
                    except Exception as e:
                        log_info(f"generate_json_structure :: PPTX Generate Json failed for Slide:{id_sld}, file:{file_id}"
                                 f"ERROR:{str(e)}", None)
            except Exception as e:
                log_info(f"generate_json_structure :: PPTX Generate Json failed for File:{file_id}"
                         f"ERROR:{str(e)}", None)
        return base_json

    def write_json_file(self, transformed_obj):
        out_file_name = config.PPTX_FILE_PREFIX + self.file_name_without_ext + '.json'
        return common_obj.write_json_file(out_file_name=out_file_name, transformed_obj=transformed_obj)

    def distribute_over_runs(self, iterable_obj, trans_para):
        common_obj.distribute_over_runs(iterable_obj=iterable_obj, trans_para=trans_para)

    def translate_pptx_file(self, document, trans_map):
        file_id = self.file_name_without_ext

        if config.PPTX_PARAGRAPH_GEN and config.PPTX_PARAGRAPH_TRANS:
            try:
                for id_sld, slide in enumerate(document.slides):
                    # NEW LOGIC TO GET PARA FOR PPTX
                    try:
                        para_lis = slide._element.xpath('.//a:p')
                        for idx, para_obj in enumerate(para_lis):
                            try:
                                para = _Paragraph(para_obj, slide)
                                runs = common_obj.get_runs(para, para_obj=True, file_type=config.TYPE_PPTX)
                                para_id = common_obj.generate_id(file_id=file_id,
                                                                 slide=str(id_sld),
                                                                 para=str(idx))
                                if para_id in trans_map:
                                    self.distribute_over_runs(runs, trans_para=trans_map[para_id])
                                else:
                                    raise FileErrors(f"translate_pptx_file:", f"PARA ID :{para_id} not found in fetch content")
                            except Exception as e:
                                log_info(f"translate_pptx_file :: Translation failed for Paragraph:{idx}, slide:{id_sld}, file:{file_id}"
                                         f"ERROR:{str(e)}", None)
                    except Exception as e:
                        log_info(f"translate_pptx_file :: Translation failed for Slide:{id_sld}, file:{file_id}"
                                 f"ERROR:{str(e)}", None)
            except Exception as e:
                log_info(f"translate_pptx_file :: Translation failed for File:{file_id}"
                         f"ERROR:{str(e)}", None)


        return document

    def write_pptx_file(self, document):
        file_name = str(uuid.uuid4()) + '.pptx'
        file_out_path = common_obj.input_path(file_name)
        document.save(file_out_path)
        return file_name
